import inspect
from io import BytesIO

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pytest
from pptx import Presentation

from pptx_replace import replace_picture, replace_table, replace_text


@pytest.fixture
def prs():
    return Presentation("tests/templates/test_template.pptx")


def test_replace_text(prs) -> None:
    replace_text(prs, "{Main title}", "this is main report title")
    slide = prs.slides[1]
    replace_text(slide, "{title}", "This is a title")
    replace_text(slide, "{content}", "a quick brown fox jumps over the lazy dog\n" * 5)

    prs.save(f"/tmp/{inspect.stack()[0][3]}.pptx")


def test_replace_picture(prs) -> None:
    # generate fig
    fig_file = BytesIO()
    plt.plot([1, 2, 3, 4])
    fig = plt.gcf()
    fig.savefig(fig_file, format="png")
    fig_file.seek(0)

    # replace picture

    # replace the first picture in slide 0
    replace_picture(prs.slides[0], fig_file, auto_reshape=True)

    # replace the first picture in slide 1
    plt.bar(
        "picture in slide is replaced in order of top to bottom".split("\n"), range(11)
    )
    fig = plt.gcf()
    replace_picture(prs.slides[1], fig, auto_reshape=False, order="l2r")

    # replace the second picture in slide 1 with out auto_reshape
    replace_picture(prs.slides[1], fig, pic_number=1, auto_reshape=True, order="l2r")

    prs.save(f"/tmp/{inspect.stack()[0][3]}.pptx")


def test_replace_altair_chart(prs) -> None:
    # open pptx file
    prs = Presentation("tests/templates/test_template.pptx")
    # generate altair chart
    import altair as alt
    import pandas as pd

    source = pd.DataFrame(
        {
            "a": ["A", "B", "C", "D", "E", "F", "G", "H", "I"],
            "b": [28, 55, 43, 91, 81, 53, 19, 87, 52],
        }
    )

    c1 = alt.Chart(source).mark_bar().encode(x="a", y="b")

    c2 = c1.mark_line()

    replace_picture(prs.slides[1], c1, auto_reshape=False, order="l2r")

    # replace the second picture in slide 1 with out auto_reshape
    replace_picture(
        prs.slides[1], c1 + c2, pic_number=1, auto_reshape=True, order="l2r"
    )
    prs.save(f"/tmp/{inspect.stack()[0][3]}.pptx")


def test_replace_table(prs) -> None:
    # open pptx file
    prs = Presentation("tests/templates/test_template.pptx")
    slide = prs.slides[3]
    # generate table
    df = pd.DataFrame(np.random.rand(6, 10))
    df_styled = df.style.background_gradient()

    replace_table(slide, df_styled)
    prs.save(f"/tmp/{inspect.stack()[0][3]}.pptx")
