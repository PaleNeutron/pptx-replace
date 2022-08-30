from typing import Literal

from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide


def get_shape(
    slide: Slide,
    pic_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
    shape_type: Literal["picture", "table"] = "picture",
) -> BaseShape:
    if shape_type == "picture":
        _ele = [shape for shape in slide.shapes if isinstance(shape, Picture)]
    elif shape_type == "table":
        _ele = [shape for shape in slide.shapes if shape.has_table]

    if order == "t2b":
        ordered_shapes = sorted(
            _ele,
            key=lambda x: (x.top, x.left),  # type: ignore
        )
    elif order == "l2r":
        ordered_shapes = sorted(
            _ele,
            key=lambda x: (x.left, x.top),  # type: ignore
        )
    else:
        raise ValueError("order must be t2b or l2r")

    shape = ordered_shapes[pic_number]
    return shape


def get_slide_from_shape(shape: BaseShape) -> Slide:
    for i in range(100):
        shape = shape._parent
        if isinstance(shape, Slide):
            return shape
    else:
        raise ValueError("shape is not in a slide")
