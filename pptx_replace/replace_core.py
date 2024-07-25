import html
import re
from io import BytesIO, IOBase
from typing import BinaryIO, List, Literal, Optional, Union

from python_docx_replace import Paragraph as DocParagraph

from pptx_replace.text import set_frame_text

try:
    import altair as alt

    HAS_ALT = True
except ImportError:
    HAS_ALT = False

import matplotlib.pyplot as plt
import pandas as pd
from pandas.io.formats.style import Styler
from pptx.presentation import Presentation as PrsCls
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from .utils import get_shape, get_slide_from_shape

__all__ = [
    "replace_text",
    "replace_picture",
    "replace_table",
]

def get_all_paragraphs(slide: Slide):
    # get from shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield from shape.text_frame.paragraphs
    # get from tables
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield from cell.text_frame.paragraphs

def remove_row(table, row):
    tbl = table._tbl
    tr = row._tr
    tbl.remove(tr)

def remove_column(table, column):
    col_idx = table._tbl.tblGrid.index(column._gridCol)
    for tr in table._tbl.tr_lst:
        tr.remove(tr.tc_lst[col_idx])
    table._tbl.tblGrid.remove(column._gridCol)

def replace_text(ppt: Union[PrsCls, Slide], search_pattern: str, repl: Optional[str]=None) -> None:
    """search and replace text in PowerPoint while preserving formatting

    Args:
        ppt: input pptx file, Presentation object / slide object
        search_pattern: pattern to search
        repl: replacement string

    Returns:
        ppt object
    """
    # Useful Links ;)
    # https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    # https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    if repl is None:
        repl = ""
    if isinstance(ppt, PrsCls):
        for slide in ppt.slides:
            _replace_text_in_slide(slide, search_pattern, repl)
    elif isinstance(ppt, Slide):
        _replace_text_in_slide(ppt, search_pattern, repl)


def _replace_text_in_slide(slide: Slide, search_str: str, repl: str) -> Slide:
    """Replace text within a page of ppt and keep the format of the corresponding text.

    Args:
        slide: slide object
        search_pattern: text to search
        repl: replacement string

    Returns:
        modified slide object
    """
    paragraphs = get_all_paragraphs(slide)
    for p in paragraphs:
        dp = DocParagraph(p)
        dp.replace_key(search_str, repl)

def replace_picture(
    slide: Slide,
    fig: Union[
        bytes,
        str,
        plt.Figure,
        BinaryIO,
        "alt.Chart",
    ],
    pic_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
    auto_reshape: bool = True,
) -> BaseShape:
    """Replace pictures in PPT in a page

    Args:
        slide: PPT object
        fig: figure to fill in the slide, could be figure object / file like object / bytes / file path
        pic_number: The number of the replaced picture in the page, sorted by picture position, the default is 0
        order: Replace the order of pictures, t2b means top to bottom, l2r means left to right, the default is t2b
    """
    shape = get_shape(
        slide,
        pic_number,
        order,
        shape_type="picture",
    )
    new_shape = replace_shape_with_picture(shape, fig, auto_reshape)
    return new_shape


# def copy_row_insert_after(
# row:_RowCollection, copy_idx:int=-1, insert_idx:int=-1,
#  init_cell_func: Optional[Callable[[_Cell], None]] = None):
#     '''
#     Duplicates target row to keep
#     formatting and resets it's cells text_frames
#     (e.g. ``row = table.rows.copy_row_insert_after(0,1)``, which copies
#     the first row and inserts after the second row as new third row).
#     Returns new |_Row| instance.
#     '''
#     new_row = copy.deepcopy(row._tbl.tr_lst[copy_idx])  # copies idx row
#     element. Note: tr_lst idx is != _tbl idx.

#     for tc in new_row.tc_lst:
#         cell = _Cell(tc, new_row.tc_lst)
#         if init_cell_func:
#             init_cell_func(cell)

#     #_tbl[0] xml sets up the table and relationships <a:tblPr>: try table.rows.debug_tbl_idx(0)
#         #https://python-pptx.readthedocs.io/en/latest/dev/analysis/tbl-table.html?highlight=a%3AtblPr#xml-semantics
#     #_tbl[1] xml sets up the columns <a:tblGrid>: try table.rows.debug_tbl_idx(1)
#     #_tbl[2] xml is the first row <a:tr>: try table.rows.debug_tbl_idx(2)

#     self._tbl.insert(insert_idx, new_row) #rows begin starting idx 2. Need to read _tbl[0], _tbl[1] xml.

#     return _Row(new_row, self)


def replace_table(
    slide: Slide,
    data: Union[
        pd.DataFrame,
        Styler,
        List[List[Union[float, str]]],
    ],
    shape_number: int = 0,
    order: Literal["t2b", "l2r"] = "t2b",
    font = None,
) -> BaseShape:
    """Replace table in PPT in a page

    Args:
        slide: PPT object
        data: dataframe to fill in the slide, could be dataframe object / Styler object / list of list
        shape_number: The number of the replaced table in the page, sorted by table position, the default is 0
        order: Replace the order of tables, t2b means top to bottom, l2r means left to right, the default is t2b
    """

    shape = get_shape(
        slide,
        shape_number,
        order,
        shape_type="table",
    )
    if isinstance(data, List):
        df = pd.DataFrame(data)

    elif isinstance(data, Styler):
        data._compute()  # type: ignore
        pandas_styles = data._translate(False, False)  # type: ignore
        df = data.data  # type: ignore
    elif isinstance(data, pd.DataFrame):
        df = data
    else:
        raise ValueError(f"{type(data)} {repr(data)} is not supported")
    
    # x, y, cx, cy = (
    #     shape.left,
    #     shape.top,
    #     shape.width,
    #     shape.height,
    # )
    if font is None:
        font = shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font
    # t = shape.table
    rn, cn = df.shape
    # shape = slide.shapes.add_table(rn + 1, cn + 1, x, y, cx, cy)
    # alt table rows and columns to meet dataframe shape
    if rn + 1 > len(shape.table.rows):
        for _ in range(rn + 1 - len(shape.table.rows)):
            shape.table.rows.add()
    elif rn + 1 < len(shape.table.rows):
        for _ in range(len(shape.table.rows) - rn - 1):
            remove_row(shape.table, shape.table.rows[-1])
    if cn + 1 > len(shape.table.columns):
        for _ in range(cn + 1 - len(shape.table.columns)):
            shape.table.columns.add()
    elif cn + 1 < len(shape.table.columns):
        for _ in range(len(shape.table.columns) - cn - 1):
            remove_column(shape.table, shape.table.columns[-1])

    # add headers
    for c in range(cn):
        if isinstance(data, pd.DataFrame):
            shape.table.cell(0, c + 1).text = str(df.columns[c])
        else:
            shape.table.cell(0, c + 1).text = html.unescape(
            pandas_styles["head"][0][c]["display_value"]
        )
    # add index
    for r in range(rn):
        if isinstance(data, pd.DataFrame):
            shape.table.cell(r + 1, 0).text = str(df.index[r])
        else:
            shape.table.cell(r + 1, 0).text = html.unescape(
                pandas_styles["body"][r][0]["display_value"]
            )
    # add body
    for r in range(rn):
        for c in range(cn):
            # tc = copy.deepcopy(shape.table.cell(-1, -1)._tc)
            # new_shape.table.cell(r+1, c)._tc = tc
            if isinstance(data, pd.DataFrame):
                shape.table.cell(r + 1, c + 1).text = str(df.iloc[r, c])
            else:
                shape.table.cell(r + 1, c + 1).text = html.unescape(
                    pandas_styles["body"][r][c]["display_value"]
                )
    # set font
    for r in range(rn + 1):
        for c in range(cn):
            for p in shape.table.cell(r, c).text_frame.paragraphs:
                for run in p.runs:
                    run.font.bold = font.bold
                    run.font.italic = font.italic
                    run.font.size = font.size
                    run.font.name = font.name
                    run.font.underline = font.underline
    old_shape = shape._element
    new_element = shape._element
    old_shape.addnext(new_element)
    old_shape.getparent().remove(old_shape)
    return shape

def replace_table_cells(
    shape: BaseShape,
    data: Union[
        pd.DataFrame,
        List[List[Union[float, str]]],
    ],
    replace_headers: bool = True,
    replace_index: bool = True,
) -> BaseShape:
    if isinstance(data, List):
        df = pd.DataFrame(data)
    else:
        df = data
    min_col = min(len(df.columns), len(shape.table.columns) -1 )
    min_row = min(len(df.index), len(shape.table.rows) -1 )
    if replace_headers:
        for c in range(1, min_col):
            set_frame_text(shape.table.cell(0, c).text_frame, str(df.columns[c]))
    if replace_index:
        for r in range(1, min_row):
            set_frame_text(shape.table.cell(r, 0).text_frame, str(df.index[r]))
    
    for r in range(1, min_row):
        for c in range(1, min_col):
            set_frame_text(shape.table.cell(r+1, c+1).text_frame, str(df.iloc[r, c]))
    return shape

def replace_shape_with_picture(
    shape: BaseShape,
    fig: Union[
        bytes,
        str,
        plt.Figure,
        BinaryIO,
        "alt.Chart",
    ],
    auto_reshape: bool = True,
    resize: int = 1,
) -> BaseShape:
    # prepare figure
    if isinstance(fig, str) or isinstance(fig, IOBase):
        figio = fig
    elif isinstance(fig, bytes):
        figio = BytesIO(fig)
    elif isinstance(fig, plt.Figure):
        figio = BytesIO()
        if auto_reshape:
            fig.set_size_inches(
                shape.width.inches * resize, shape.height.inches * resize
            )
        if fig.get_constrained_layout():
            fig.savefig(figio, format="png")
        else:
            fig.savefig(figio, format="png", bbox_inches="tight")
    elif HAS_ALT and isinstance(fig, alt.VegaLiteSchema):
        figio = BytesIO()
        if auto_reshape:
            fig = fig.properties(
                width=shape.width.pt * resize, height=shape.height.pt * resize
            )
        fig.save(figio, format="png")
    else:
        raise ValueError(f"{type(fig)} {repr(fig)} is not supported")

    # replace picture
    slide = get_slide_from_shape(shape)
    new_shape = slide.shapes.add_picture(
        figio,
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    old_pic = shape._element
    new_pic = new_shape._element
    old_pic.addnext(new_pic)
    old_pic.getparent().remove(old_pic)
    return new_shape