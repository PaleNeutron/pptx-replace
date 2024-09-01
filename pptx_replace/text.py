from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.text.text import Font, TextFrame


def delete_paragraph(paragraph):
    p = paragraph._p
    parent_element = p.getparent()
    parent_element.remove(p)


def set_frame_text(text_frame: TextFrame, text: str, font: Font = None):
    """set text and keep font style"""
    replaced = False
    if text_frame.paragraphs:
        for paragraph in text_frame.paragraphs:
            # use style of the first run
            if paragraph.runs:
                if not replaced:
                    paragraph.runs[0].text = text
                    # remove text in other runs
                    replaced = True
                    for run in paragraph.runs[1:]:
                        run.text = ""
                else:
                    delete_paragraph(paragraph)
            else:
                paragraph.text = text
    else:
        text_frame.text = text

    # actually after set text, there should be only one paragraph and one run
    if font is not None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = font.bold
                run.font.italic = font.italic
                run.font.size = font.size
                run.font.name = font.name
                run.font.underline = font.underline
                # run.font.color.rgb = font.color.rgb
    return text_frame


def table_alignment(table, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER):
    """set alignment of cell in table"""
    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = alignment
    return table
